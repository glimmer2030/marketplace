#!/usr/bin/env python3
"""
Generate Weekly Report PowerPoint from markdown + template.

Usage:
    python3 generate_pptx.py <week_number> <markdown_file> <output_pptx> [options]

Options:
    --template PATH     Path to .potx template file
    --date-range RANGE  Date range string (e.g. "2/10~2/21")

Examples:
    python3 generate_pptx.py 7-8 week-7-8.md week-7-8.pptx --template template.potx --date-range "2/10~2/21"
    python3 generate_pptx.py 3 week-3.md week-3.pptx  # fallback: no template

Requirements:
    pip install defusedxml
    pptx skill installed at ~/.claude/skills/pptx/ (for template mode)
"""

import argparse
import os
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

# Locate pptx skill scripts
PPTX_SKILL_PATHS = [
    os.path.expanduser("~/.claude/skills/pptx/scripts"),
    os.path.expanduser("~/.agents/skills/pptx/scripts"),
]


def find_pptx_skill():
    """Find the pptx skill scripts directory."""
    for p in PPTX_SKILL_PATHS:
        if os.path.isdir(p):
            return p
    return None


def find_venv_python():
    """Find a Python with defusedxml available."""
    # Check if current python has defusedxml
    try:
        import defusedxml  # noqa: F401
        return sys.executable
    except ImportError:
        pass

    # Check common venv locations
    candidates = [
        os.path.join(os.getcwd(), ".venv/bin/python3"),
        os.path.expanduser("~/projects/weekly-report/.venv/bin/python3"),
    ]
    for c in candidates:
        if os.path.isfile(c):
            return c

    return sys.executable


# ---------------------------------------------------------------------------
# Markdown parser
# ---------------------------------------------------------------------------

def parse_markdown(md_file):
    """Parse markdown file and extract tasks with their sections."""
    with open(md_file, "r", encoding="utf-8") as f:
        content = f.read()

    tasks = []
    current_task = None
    current_section = None
    current_content = []
    in_format_spec = False

    for line in content.split("\n"):
        if line.strip() == "---":
            if in_format_spec:
                in_format_spec = False
                continue
            elif current_task:
                if current_section:
                    current_task["sections"][current_section] = current_content
                tasks.append(current_task)
                current_task = None
                current_section = None
                current_content = []
            continue

        if in_format_spec or line.startswith("**格式規範"):
            in_format_spec = True
            continue

        if line.startswith("## ") and ("Task" in line or "任務" in line or ":" in line):
            if current_task and current_section:
                current_task["sections"][current_section] = current_content
            if current_task:
                tasks.append(current_task)

            task_name = line[3:].strip()
            # Strip "Task N: " prefix for display
            short_name = re.sub(r"^Task\s+\d+:\s*", "", task_name)
            current_task = {"name": task_name, "short_name": short_name, "sections": {}}
            current_section = None
            current_content = []

        elif line.startswith("### "):
            if current_section and current_task:
                current_task["sections"][current_section] = current_content
            current_section = line[4:].strip()
            current_content = []

        elif line.startswith("  ") and line.strip() and current_section:
            current_content.append(line.strip())

    if current_task and current_section:
        current_task["sections"][current_section] = current_content
    if current_task:
        tasks.append(current_task)

    return tasks


def infer_status(task):
    """Infer task status from content. Has real solutions → Close, otherwise In progress."""
    solutions = task["sections"].get("解決方案", [])
    if not solutions or all(s.strip() == "無" for s in solutions):
        return "In progress"
    return "Close"


# ---------------------------------------------------------------------------
# Template-based generation
# ---------------------------------------------------------------------------

def convert_potx_to_pptx(potx_path, pptx_path):
    """Convert .potx to .pptx by changing content type in zip."""
    with zipfile.ZipFile(potx_path, "r") as zin, zipfile.ZipFile(pptx_path, "w") as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = data.replace(
                    b"presentationml.template.main+xml",
                    b"presentationml.presentation.main+xml",
                )
            zout.writestr(item, data)


def run_skill_script(skill_dir, script_name, args, python=None):
    """Run a pptx skill script."""
    script_path = os.path.join(skill_dir, script_name)
    if not os.path.isfile(script_path):
        # Try in office/ subdirectory
        script_path = os.path.join(skill_dir, "office", script_name)
    if not os.path.isfile(script_path):
        print(f"Error: Script not found: {script_name}", file=sys.stderr)
        sys.exit(1)

    python = python or find_venv_python()
    env = os.environ.copy()
    # Add office dir to PYTHONPATH for pack.py's validators import
    office_dir = os.path.join(skill_dir, "office")
    env["PYTHONPATH"] = office_dir + ":" + env.get("PYTHONPATH", "")

    result = subprocess.run(
        [python, script_path] + args,
        capture_output=True, text=True, env=env,
    )
    if result.returncode != 0:
        error_msg = result.stderr or result.stdout or "(no output)"
        print(f"Script {script_name} failed:\n{error_msg}", file=sys.stderr)
        sys.exit(1)
    return result.stdout


def _read_sld_id_list(pres_path):
    """Read existing sldIdLst entries from presentation.xml.
    Returns list of (id, r:id) tuples in order."""
    with open(pres_path, "r", encoding="utf-8") as f:
        content = f.read()
    entries = re.findall(r'<p:sldId[^>]*\bid="(\d+)"[^>]*r:id="(rId\d+)"', content)
    if not entries:
        # Try reversed attribute order
        entries = re.findall(r'<p:sldId[^>]*r:id="(rId\d+)"[^>]*\bid="(\d+)"', content)
        entries = [(sid, rid) for rid, sid in entries]
    return entries


def _read_rels_rid_to_slide(unpacked):
    """Read presentation.xml.rels to map rId → slide filename."""
    rels_path = os.path.join(unpacked, "ppt", "_rels", "presentation.xml.rels")
    with open(rels_path, "r", encoding="utf-8") as f:
        content = f.read()
    rid_map = {}
    for m in re.finditer(r'Id="(rId\d+)"[^>]*Target="slides/([^"]+)"', content):
        rid_map[m.group(1)] = m.group(2)
    # Also try reversed attribute order
    for m in re.finditer(r'Target="slides/([^"]+)"[^>]*Id="(rId\d+)"', content):
        rid_map[m.group(2)] = m.group(1)
    return rid_map


def generate_with_template(tasks, template_path, output_path, date_range):
    """Generate PPTX using template."""
    from pptx_templates import sub_item_xml, BLANK_XML, xml_escape

    skill_dir = find_pptx_skill()
    if not skill_dir:
        print("Error: pptx skill not installed. Install with: npx skills add anthropics/skills@pptx -g -y", file=sys.stderr)
        sys.exit(1)

    python = find_venv_python()
    work_dir = tempfile.mkdtemp(prefix="weekly-report-")
    converted = os.path.join(work_dir, "converted.pptx")

    try:
        # Step 1: Convert .potx → .pptx
        print("Converting template...")
        convert_potx_to_pptx(template_path, converted)

        # Step 2: Unpack
        unpacked = os.path.join(work_dir, "unpacked")
        print("Unpacking template...")
        run_skill_script(skill_dir, "unpack.py", [converted, unpacked], python)

        # Step 2b: Read original slide structure
        pres_path = os.path.join(unpacked, "ppt", "presentation.xml")
        orig_entries = _read_sld_id_list(pres_path)
        rid_to_slide = _read_rels_rid_to_slide(unpacked)

        # Identify slide roles by position:
        # [0]=cover, [1]=summary, [2]=task template, [-2]=thanks, [-1]=appendix
        # The task template slide (index 2) will be duplicated for extra tasks
        if len(orig_entries) < 4:
            print("Error: Template must have at least 4 slides (cover, summary, task, thanks)", file=sys.stderr)
            sys.exit(1)

        cover_entry = orig_entries[0]
        summary_entry = orig_entries[1]
        task_template_entry = orig_entries[2]
        tail_entries = orig_entries[3:]  # thanks, appendix, etc.

        task_template_slide = rid_to_slide.get(task_template_entry[1], "slide3.xml")

        # Step 3: Duplicate task template slide for extra tasks
        num_tasks = len(tasks)
        new_rids = []  # rIds from add_slide.py for duplicated slides
        if num_tasks > 1:
            print(f"Duplicating {task_template_slide} for {num_tasks} tasks...")
            for i in range(num_tasks - 1):
                output = run_skill_script(skill_dir, "add_slide.py", [unpacked, task_template_slide], python)
                # Parse rId from output like: <p:sldId id="261" r:id="rId53"/>
                match = re.search(r'r:id="(rId\d+)"', output)
                if match:
                    new_rids.append(match.group(1))

        # Step 4: Build and write correct sldIdLst
        print("Reordering slides...")
        with open(pres_path, "r", encoding="utf-8") as f:
            pres_content = f.read()

        # Assign unique IDs: start from max existing + 1
        max_id = max(int(sid) for sid, _ in orig_entries)

        slide_entries = []
        slide_entries.append(f'    <p:sldId id="{cover_entry[0]}" r:id="{cover_entry[1]}"/>')
        slide_entries.append(f'    <p:sldId id="{summary_entry[0]}" r:id="{summary_entry[1]}"/>')
        slide_entries.append(f'    <p:sldId id="{task_template_entry[0]}" r:id="{task_template_entry[1]}"/>')

        for rid in new_rids:
            max_id += 1
            slide_entries.append(f'    <p:sldId id="{max_id}" r:id="{rid}"/>')

        for sid, rid in tail_entries:
            slide_entries.append(f'    <p:sldId id="{sid}" r:id="{rid}"/>')

        new_sld_list = "  <p:sldIdLst>\n" + "\n".join(slide_entries) + "\n  </p:sldIdLst>"

        # Replace existing sldIdLst
        pres_content = re.sub(
            r"<p:sldIdLst>.*?</p:sldIdLst>",
            new_sld_list,
            pres_content,
            flags=re.DOTALL,
        )
        with open(pres_path, "w", encoding="utf-8") as f:
            f.write(pres_content)

        # Step 5: Fill cover date
        print("Filling cover slide...")
        cover_slide = rid_to_slide.get(cover_entry[1], "slide1.xml")
        slide1_path = os.path.join(unpacked, "ppt", "slides", cover_slide)
        if os.path.isfile(slide1_path):
            with open(slide1_path, "r", encoding="utf-8") as f:
                slide1 = f.read()
            # Update dates if date_range provided
            if date_range:
                # Try to find and replace existing date patterns
                parts = re.split(r"[~～\-]", date_range)
                if len(parts) == 2:
                    start, end = parts[0].strip(), parts[1].strip()
                    # Replace dates: first occurrence → start, rest → end
                    count = [0]
                    def _replace_dot_date(m):
                        count[0] += 1
                        return f"2026.{start.replace('/', '.')}" if count[0] == 1 else f"2026.{end.replace('/', '.')}"
                    slide1 = re.sub(r"20\d{2}\.\d{1,2}\.\d{1,2}", _replace_dot_date, slide1)
                    slide1 = re.sub(r"20\d{2}/\d+/\d+", f"2026/{end}", slide1)
            with open(slide1_path, "w", encoding="utf-8") as f:
                f.write(slide1)

        # Step 6: Fill summary table
        print("Filling summary table...")
        summary_slide = rid_to_slide.get(summary_entry[1], "slide2.xml")
        slide2_path = os.path.join(unpacked, "ppt", "slides", summary_slide)
        fill_summary_table(slide2_path, tasks, date_range)

        # Step 7: Fill task detail slides
        print("Filling task slides...")
        # Build task slide list: first is the template slide, rest are from new_rids
        # Re-read rels since add_slide.py updated it
        updated_rid_to_slide = _read_rels_rid_to_slide(unpacked)
        task_slide_files = [task_template_slide]
        for rid in new_rids:
            slide_name = updated_rid_to_slide.get(rid)
            if slide_name:
                task_slide_files.append(slide_name)

        slides_dir = os.path.join(unpacked, "ppt", "slides")

        for i, task in enumerate(tasks):
            if i < len(task_slide_files):
                slide_path = os.path.join(slides_dir, task_slide_files[i])
                status = infer_status(task)
                fill_task_slide(slide_path, task, date_range, status)

        # Step 8: Fix Content_Types for GIF if needed
        ct_path = os.path.join(unpacked, "[Content_Types].xml")
        with open(ct_path, "r", encoding="utf-8") as f:
            ct_content = f.read()
        if 'Extension="gif"' not in ct_content:
            ct_content = ct_content.replace(
                'Extension="jpeg" ContentType="image/jpeg"/>',
                'Extension="jpeg" ContentType="image/jpeg"/>\n  <Default Extension="gif" ContentType="image/gif"/>',
            )
            with open(ct_path, "w", encoding="utf-8") as f:
                f.write(ct_content)

        # Step 9: Ensure validators __init__.py exists
        validators_dir = os.path.join(skill_dir, "office", "validators")
        init_path = os.path.join(validators_dir, "__init__.py")
        if os.path.isdir(validators_dir) and not os.path.isfile(init_path):
            with open(init_path, "w") as f:
                f.write("from .docx import DOCXSchemaValidator\n")
                f.write("from .pptx import PPTXSchemaValidator\n")
                f.write("from .redlining import RedliningValidator\n")

        # Step 10: Clean and pack
        print("Cleaning and packing...")
        run_skill_script(skill_dir, "clean.py", [unpacked], python)
        run_skill_script(skill_dir, "pack.py", [unpacked, output_path, "--original", converted], python)

        print(f"Generated: {output_path}")
        print(f"Slides: {2 + num_tasks + 2} (cover + summary + {num_tasks} tasks + thanks + appendix)")

    finally:
        shutil.rmtree(work_dir, ignore_errors=True)


def fill_summary_table(slide2_path, tasks, date_range):
    """Fill the summary table in slide 2 using defusedxml."""
    import defusedxml.minidom as minidom

    with open(slide2_path, "r", encoding="utf-8") as f:
        doc = minidom.parseString(f.read())

    # Update title date
    if date_range:
        for r in doc.getElementsByTagName("a:r"):
            t_nodes = r.getElementsByTagName("a:t")
            if t_nodes and t_nodes[0].firstChild:
                text = t_nodes[0].firstChild.nodeValue
                if text.startswith("Summary "):
                    t_nodes[0].firstChild.nodeValue = f"Summary {date_range.replace('~', ' ~ ')}"

    # Fill table rows
    rows = doc.getElementsByTagName("a:tr")
    parts = re.split(r"[~～\-]", date_range) if date_range else ["", ""]
    start_date = parts[0].strip() if len(parts) > 0 else ""
    end_date = parts[1].strip() if len(parts) > 1 else ""

    for task_idx, task in enumerate(tasks):
        if task_idx + 1 >= len(rows):
            break
        row = rows[task_idx + 1]  # Skip header row
        cells = row.getElementsByTagName("a:tc")

        status = infer_status(task)
        # Split short_name into project and description for table
        short = task["short_name"]
        # Try to split on Chinese/English boundary or common patterns
        project = short.split(" ")[0] if " " in short else short
        desc = " ".join(short.split(" ")[1:]) if " " in short else ""

        values = [
            str(task_idx + 1),  # No
            project,            # Project
            desc,               # Task description
            start_date,         # Start date
            end_date,           # Due date
            end_date if status == "Close" else "",  # End date
            status,             # Status
            "William",          # Owner
            "",                 # Points
        ]

        for col_idx, value in enumerate(values):
            if col_idx < len(cells):
                _set_cell_text(doc, cells[col_idx], value)

    with open(slide2_path, "w", encoding="utf-8") as f:
        f.write(doc.toxml())


def _set_cell_text(doc, cell, text):
    """Set text in a table cell, preserving formatting."""
    tx_body = cell.getElementsByTagName("a:txBody")[0]
    p_elements = tx_body.getElementsByTagName("a:p")
    if not p_elements:
        return
    p = p_elements[0]

    r_elements = p.getElementsByTagName("a:r")
    if r_elements:
        t = r_elements[0].getElementsByTagName("a:t")[0]
        while t.firstChild:
            t.removeChild(t.firstChild)
        t.appendChild(doc.createTextNode(text))
    else:
        end_rpr = p.getElementsByTagName("a:endParaRPr")
        if end_rpr:
            new_rpr = doc.createElement("a:rPr")
            for attr in end_rpr[0].attributes.keys():
                new_rpr.setAttribute(attr, end_rpr[0].getAttribute(attr))
            for child in list(end_rpr[0].childNodes):
                new_rpr.appendChild(child.cloneNode(True))

            r = doc.createElement("a:r")
            r.appendChild(new_rpr)
            t = doc.createElement("a:t")
            t.appendChild(doc.createTextNode(text))
            r.appendChild(t)
            p.insertBefore(r, end_rpr[0])


def fill_task_slide(slide_path, task, date_range, status):
    """Fill a task detail slide with content."""
    from pptx_templates import sub_item_xml, BLANK_XML, xml_escape

    with open(slide_path, "r", encoding="utf-8") as f:
        content = f.read()

    # Replace title
    content = re.sub(
        r"(<a:t>)1/6-1/12, (</a:t>)",
        rf"\g<1>{xml_escape(date_range)}, \2",
        content,
    )
    content = re.sub(
        r"(<a:t>)測試模型量化效能(</a:t>)",
        rf"\g<1>{xml_escape(task['short_name'])}\2",
        content,
    )

    # Replace status
    status_text = f"({status})"
    content = content.replace(
        "<a:t>(Close/In progress/Pending)</a:t>",
        f"<a:t>{status_text}</a:t>",
    )

    # Section mapping: markdown key → (anchor text, items)
    section_map = [
        ("Objective", "<a:t>目標</a:t>"),
        ("本週進度", "<a:t>(Close)</a:t>" if status == "Close" else f"<a:t>({status})</a:t>"),
        ("困難", "<a:t>困難</a:t>"),
        ("解決方案", "<a:t>方向 </a:t>"),
    ]

    for section_key, anchor in section_map:
        items = task["sections"].get(section_key, [])
        if not items:
            continue

        # Find the anchor and locate end of its </a:p>
        anchor_idx = content.find(anchor)
        if anchor_idx == -1:
            continue

        end_p_idx = content.index("</a:p>", anchor_idx) + len("</a:p>")

        # Build sub-items XML
        items_xml = "\n".join(sub_item_xml(item) for item in items)

        # Insert after the anchor paragraph
        content = content[:end_p_idx] + "\n" + items_xml + "\n" + content[end_p_idx:]

    with open(slide_path, "w", encoding="utf-8") as f:
        f.write(content)


# ---------------------------------------------------------------------------
# Fallback: create from scratch (no template)
# ---------------------------------------------------------------------------

def generate_from_scratch(tasks, week_num, output_path):
    """Fallback: generate PPTX from scratch using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Pt, Inches
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except ImportError:
        print("Error: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    FONT_CHINESE = "微軟正黑體"
    SECTION_HEADER_SIZE = Pt(18)
    CONTENT_SIZE = Pt(14)

    def has_chinese(text):
        return bool(re.search(r"[\u4e00-\u9fff]", text))

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = title_slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = f"Week {week_num} 週報"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.name = FONT_CHINESE
    title_para.font.color.rgb = RGBColor(0, 0, 0)

    section_mapping = {
        "Objective": "目標",
        "本週進度": "進度",
        "困難": "困難",
        "解決方案": "解決方案/方向",
    }

    for task in tasks:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        tf = title_box.text_frame
        tf.text = task["name"]
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.name = FONT_CHINESE if has_chinese(task["name"]) else "Calibri"
        p.font.color.rgb = RGBColor(0, 0, 0)

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(6))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        first_section = True
        for key, display in section_mapping.items():
            if key not in task["sections"]:
                continue
            if not first_section:
                text_frame.add_paragraph().text = ""
            first_section = False

            p = text_frame.add_paragraph()
            p.text = display
            p.font.size = SECTION_HEADER_SIZE
            p.font.bold = True
            p.font.name = FONT_CHINESE
            p.font.color.rgb = RGBColor(0, 0, 0)

            for item in task["sections"][key]:
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = CONTENT_SIZE
                p.font.name = FONT_CHINESE if has_chinese(item) else "Calibri"
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.level = 1

    prs.save(output_path)
    print(f"Generated (from scratch): {output_path}")
    print(f"Slides: {1 + len(tasks)} (1 title + {len(tasks)} task slides)")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="Generate Weekly Report PowerPoint")
    parser.add_argument("week_number", help="Week number (e.g. 7-8)")
    parser.add_argument("markdown_file", help="Path to week-*.md file")
    parser.add_argument("output_pptx", help="Output .pptx file path")
    parser.add_argument("--template", help="Path to .potx template file")
    parser.add_argument("--date-range", help='Date range (e.g. "2/10~2/21")')
    args = parser.parse_args()

    if not os.path.isfile(args.markdown_file):
        print(f"Error: Markdown file not found: {args.markdown_file}", file=sys.stderr)
        sys.exit(1)

    # Add scripts dir to path for pptx_templates import
    sys.path.insert(0, SCRIPT_DIR)

    tasks = parse_markdown(args.markdown_file)
    if not tasks:
        print("Error: No tasks found in markdown file", file=sys.stderr)
        sys.exit(1)

    print(f"Parsed {len(tasks)} tasks from {args.markdown_file}")

    if args.template and os.path.isfile(args.template):
        if not args.date_range:
            print("Warning: --date-range not specified, dates may not be updated", file=sys.stderr)
        generate_with_template(tasks, args.template, args.output_pptx, args.date_range or "")
    else:
        if args.template:
            print(f"Warning: Template not found: {args.template}, falling back to scratch mode", file=sys.stderr)
        generate_from_scratch(tasks, args.week_number, args.output_pptx)


if __name__ == "__main__":
    main()
